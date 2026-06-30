"""
Activity 1: Research Methodology Selection - PPTX
Topic: Detection of APT in Cloud Networks Using AI
4 slides with professional dark theme and humanized wording
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Colors
BG = RGBColor(0x0B, 0x11, 0x20)
CARD = RGBColor(0x1A, 0x23, 0x32)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LT = RGBColor(0x93, 0xC5, 0xFD)
TEAL = RGBColor(0x08, 0x91, 0xB2)
TEAL_LT = RGBColor(0x67, 0xE8, 0xF9)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LT = RGBColor(0xC4, 0xB5, 0xFD)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_LT = RGBColor(0x6E, 0xE7, 0xB7)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SEC = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def orb(slide, l, t, w, h, color, alpha=0.87):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    spPr = s._element.find(qn('p:spPr'))
    if spPr is not None:
        sf = spPr.find('.//' + qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha'))
                a.set('val', str(int((1 - alpha) * 100000)))

def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tb

def decor(slide):
    orb(slide, Inches(10.5), Inches(-0.8), Inches(3), Inches(3), BLUE)
    orb(slide, Inches(-0.8), Inches(5.5), Inches(2.2), Inches(2.2), TEAL, 0.9)

def snum(slide, n):
    txt(slide, Inches(11.8), Inches(6.95), Inches(1.2), Inches(0.35),
        f"{n:02d} / 04", sz=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: Introduction & Problem Statement
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)
decor(s1)
orb(s1, Inches(5), Inches(2.5), Inches(4), Inches(4), PURPLE, 0.93)

txt(s1, Inches(2.5), Inches(0.6), Inches(8), Inches(0.35),
    "ACTIVITY 1: RESEARCH METHODOLOGY SELECTION",
    sz=9, bold=True, color=BLUE_LT, align=PP_ALIGN.CENTER)
txt(s1, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.2),
    "Identifying the Right Research Methodology",
    sz=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, Inches(1.8), Inches(2.4), Inches(9.7), Inches(1.0),
    "Our study focuses on detecting Advanced Persistent Threats in cloud networks "
    "using artificial intelligence. Before diving into experiments, we need to select "
    "a methodology that matches the nature of our research problem.",
    sz=12, color=SEC, align=PP_ALIGN.CENTER)

# Left card - Background
rect(s1, Inches(0.5), Inches(3.7), Inches(6.0), Inches(3.4), CARD, BLUE)
txt(s1, Inches(0.8), Inches(3.9), Inches(5.4), Inches(0.4),
    "Background & Problem Statement", sz=13, bold=True, color=BLUE_LT)
txt(s1, Inches(0.8), Inches(4.4), Inches(5.4), Inches(2.6),
    "Cloud computing has become the backbone of modern enterprise infrastructure, "
    "but it has also opened new doors for sophisticated attackers. APTs are "
    "long-running campaigns that infiltrate networks and stay hidden for months. "
    "Traditional security tools rely on known signatures and simple rules, which "
    "APT actors easily bypass. We need intelligent systems that learn and adapt.",
    sz=10, color=MUTED)

# Right card - Aim & Objectives
rect(s1, Inches(6.8), Inches(3.7), Inches(6.0), Inches(3.4), CARD, TEAL)
txt(s1, Inches(7.1), Inches(3.9), Inches(5.4), Inches(0.4),
    "Research Aim & Objectives", sz=13, bold=True, color=TEAL_LT)
txt(s1, Inches(7.1), Inches(4.4), Inches(5.4), Inches(0.6),
    "Aim: Develop and evaluate an AI-driven detection framework capable of "
    "identifying multi-stage APT attacks in cloud environments.",
    sz=10, bold=True, color=WHITE)
objectives = [
    "1. Design a hybrid deep learning architecture for cloud traffic analysis",
    "2. Benchmark performance against established detection methods",
    "3. Evaluate real-time detection feasibility in dynamic cloud environments",
    "4. Provide explainable outputs that security teams can act on",
]
for i, obj in enumerate(objectives):
    txt(s1, Inches(7.3), Inches(5.1 + i * 0.42), Inches(5.2), Inches(0.38),
        obj, sz=10, color=MUTED)

snum(s1, 1)


# ============================================================
# SLIDE 2: Research Philosophies
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
decor(s2)

txt(s2, Inches(0.6), Inches(0.4), Inches(5), Inches(0.3),
    "RESEARCH PHILOSOPHY & APPROACH", sz=9, bold=True, color=PURPLE_LT)
txt(s2, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "Evaluating Research Philosophies", sz=26, bold=True, color=WHITE)
txt(s2, Inches(0.6), Inches(1.45), Inches(10), Inches(0.4),
    "Every research project rests on philosophical assumptions about how knowledge is created.",
    sz=11, color=MUTED)

philosophies = [
    ("Positivism", BLUE, BLUE_LT,
     "Knowledge from observable, measurable phenomena. Values objectivity and hypothesis testing.",
     "Perfectly suits our need to measure detection accuracy and produce reproducible results.",
     "Cannot capture strategic reasoning of threat actors or contextual nuances."),
    ("Interpretivism", TEAL, TEAL_LT,
     "Knowledge constructed through subjective human experience. Emphasises meanings and context.",
     "Could explore how analysts interpret alerts or how culture affects incident response.",
     "Our question is about AI system effectiveness, not human perceptions."),
    ("Pragmatism", PURPLE, PURPLE_LT,
     "Value of knowledge depends on practical usefulness. Supports mixing methods.",
     "Allows flexibility to combine quantitative experiments with expert feedback.",
     "Additional qualitative layer adds complexity without proportionate benefit."),
    ("Realism", GREEN, GREEN_LT,
     "Objective reality exists independently but we can only partially access it.",
     "Recognises that APT behaviour has underlying causal structures.",
     "More suited to explanatory social research than computational benchmarking."),
]

for i, (name, border, color, desc, strength, limit) in enumerate(philosophies):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.95 + row * 2.6)
    rect(s2, x, y, Inches(6.1), Inches(2.4), CARD, border)
    txt(s2, x + Inches(0.2), y + Inches(0.15), Inches(5.5), Inches(0.35),
        name, sz=12, bold=True, color=color)
    txt(s2, x + Inches(0.2), y + Inches(0.5), Inches(5.6), Inches(0.55),
        desc, sz=9, color=SEC)
    txt(s2, x + Inches(0.2), y + Inches(1.05), Inches(5.6), Inches(0.2),
        "STRENGTH:", sz=7, bold=True, color=GREEN_LT)
    txt(s2, x + Inches(1.2), y + Inches(1.05), Inches(4.6), Inches(0.4),
        strength, sz=8, color=MUTED)
    txt(s2, x + Inches(0.2), y + Inches(1.5), Inches(5.6), Inches(0.2),
        "LIMITATION:", sz=7, bold=True, color=ORANGE)
    txt(s2, x + Inches(1.3), y + Inches(1.5), Inches(4.5), Inches(0.4),
        limit, sz=8, color=MUTED)

# Bottom justification bar
rect(s2, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.0), CARD, None)

snum(s2, 2)


# ============================================================
# SLIDE 3: Methodology Selection & Justification
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
decor(s3)

txt(s3, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
    "METHODOLOGY SELECTION", sz=9, bold=True, color=GREEN_LT)
txt(s3, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "Quantitative Methodology: The Natural Choice", sz=26, bold=True, color=WHITE)
txt(s3, Inches(0.6), Inches(1.45), Inches(10), Inches(0.4),
    "Given our positivist stance, a quantitative methodology is the most appropriate path forward.",
    sz=11, color=MUTED)

# Three method comparison cards
methods = [
    ("Qualitative", PURPLE_LT, PURPLE,
     "Explores meanings through interviews and observations.",
     "Not suitable - Our research produces numerical outputs, not human perceptions."),
    ("Quantitative", BLUE_LT, BLUE,
     "Relies on measurable data, statistics, and experiments.",
     "Perfect fit - Everything is quantifiable: detection rates, accuracy, processing times."),
    ("Mixed Methods", TEAL_LT, TEAL,
     "Combines quantitative and qualitative approaches.",
     "Unnecessary - Adds complexity without proportional gain for our questions."),
]

for i, (name, color, border, desc, verdict) in enumerate(methods):
    x = Inches(0.4 + i * 4.2)
    rect(s3, x, Inches(2.0), Inches(3.9), Inches(2.0), CARD, border)
    sel_txt = "  [SELECTED]" if i == 1 else ""
    txt(s3, x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(0.35),
        name + sel_txt, sz=11, bold=True, color=color)
    txt(s3, x + Inches(0.2), Inches(2.55), Inches(3.5), Inches(0.5),
        desc, sz=9, color=SEC)
    v_color = GREEN_LT if i == 1 else ORANGE
    txt(s3, x + Inches(0.2), Inches(3.1), Inches(3.5), Inches(0.6),
        verdict, sz=9, bold=True, color=v_color)

# Left - Why Quantitative Works
rect(s3, Inches(0.4), Inches(4.2), Inches(6.2), Inches(3.0), CARD, BLUE)
txt(s3, Inches(0.7), Inches(4.4), Inches(5.5), Inches(0.35),
    "Why Quantitative Works for Us", sz=12, bold=True, color=BLUE_LT)
reasons = [
    "We train ML models on numerical datasets with millions of network flows",
    "Success is measured through concrete metrics - accuracy, precision, recall, F1",
    "We use controlled experiments with training/testing splits and cross-validation",
    "Conclusions rest on statistical significance tests, not subjective judgement",
]
for i, r in enumerate(reasons):
    txt(s3, Inches(0.9), Inches(4.9 + i * 0.5), Inches(5.5), Inches(0.45),
        f"  {r}", sz=10, color=MUTED)

# Right - Research Methods Table
rect(s3, Inches(6.9), Inches(4.2), Inches(5.8), Inches(3.0), CARD, GREEN)
txt(s3, Inches(7.2), Inches(4.4), Inches(5.2), Inches(0.35),
    "Appropriateness of Research Methods", sz=12, bold=True, color=GREEN_LT)
# Table header
rect(s3, Inches(7.0), Inches(4.85), Inches(5.5), Inches(0.35),
     RGBColor(0x12, 0x1A, 0x2A), None)
txt(s3, Inches(7.1), Inches(4.88), Inches(2.2), Inches(0.3),
    "Method", sz=8, bold=True, color=BLUE_LT)
txt(s3, Inches(9.3), Inches(4.88), Inches(3.0), Inches(0.3),
    "Role in Our Study", sz=8, bold=True, color=BLUE_LT)

table_rows = [
    ("Benchmark Datasets", "Training and baseline comparison"),
    ("Controlled Experiments", "Cloud testbed with simulated APT attacks"),
    ("Statistical Testing", "t-tests and ANOVA for comparing models"),
    ("Cross-Validation", "k-fold CV for robustness and generalisability"),
]
for i, (method, role) in enumerate(table_rows):
    y = Inches(5.25 + i * 0.45)
    txt(s3, Inches(7.1), y, Inches(2.2), Inches(0.4),
        method, sz=9, bold=True, color=WHITE)
    txt(s3, Inches(9.3), y, Inches(3.0), Inches(0.4),
        role, sz=9, color=MUTED)

snum(s3, 3)


# ============================================================
# SLIDE 4: Summary & Final Justification
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4)
decor(s4)

txt(s4, Inches(3.5), Inches(0.6), Inches(6), Inches(0.35),
    "SUMMARY", sz=9, bold=True, color=GREEN_LT, align=PP_ALIGN.CENTER)
txt(s4, Inches(2), Inches(1.1), Inches(9.3), Inches(0.7),
    "Bringing It All Together", sz=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s4, Inches(2.5), Inches(1.8), Inches(8.3), Inches(0.5),
    "Our philosophical stance and methodological choice form a coherent framework "
    "that directly supports every aspect of the research.",
    sz=12, color=SEC, align=PP_ALIGN.CENTER)

# Left - Design Overview Table
rect(s4, Inches(0.5), Inches(2.6), Inches(6.0), Inches(3.8), CARD, BLUE)
txt(s4, Inches(0.8), Inches(2.8), Inches(5.4), Inches(0.4),
    "Research Design Overview", sz=13, bold=True, color=BLUE_LT, align=PP_ALIGN.CENTER)

design = [
    ("Philosophy", "Positivism"),
    ("Approach", "Deductive"),
    ("Methodology", "Quantitative"),
    ("Primary Methods", "Experiments + Benchmark Analysis"),
    ("Data Analysis", "Statistical + ML Performance Metrics"),
]
for i, (comp, choice) in enumerate(design):
    y = Inches(3.35 + i * 0.55)
    bg = RGBColor(0x12, 0x1A, 0x2A) if i % 2 == 0 else RGBColor(0x15, 0x1E, 0x2E)
    rect(s4, Inches(0.6), y, Inches(5.7), Inches(0.48), bg, None)
    txt(s4, Inches(0.8), y + Inches(0.06), Inches(2.5), Inches(0.35),
        comp, sz=10, bold=True, color=WHITE)
    txt(s4, Inches(3.5), y + Inches(0.06), Inches(2.7), Inches(0.35),
        choice, sz=10, color=TEAL_LT)

# Right - Final Justification
rect(s4, Inches(6.8), Inches(2.6), Inches(6.0), Inches(3.8), CARD, GREEN)
txt(s4, Inches(7.1), Inches(2.8), Inches(5.4), Inches(0.4),
    "Final Justification", sz=13, bold=True, color=GREEN_LT)

justifications = [
    ("Positivism", "assumes knowledge comes from measurable observation - our AI models produce exactly that"),
    ("Deductive approach", "lets us start from established ML theory, form hypotheses, and test them"),
    ("Quantitative methodology", "provides the structured numerical framework to evaluate our system"),
    ("Experimental methods", "allow controlled, reproducible testing the community can verify"),
]
for i, (key, val) in enumerate(justifications):
    y = Inches(3.3 + i * 0.7)
    txt(s4, Inches(7.3), y, Inches(5.2), Inches(0.3),
        f"  {key}", sz=10, bold=True, color=WHITE)
    txt(s4, Inches(7.3), y + Inches(0.28), Inches(5.2), Inches(0.35),
        f"    {val}", sz=9, color=MUTED)

# Bottom note
rect(s4, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7),
     RGBColor(0x0A, 0x1A, 0x2A), BLUE)
txt(s4, Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.5),
    "Each element reinforces the others, creating a logically consistent research "
    "design suited to answering our core question with confidence.",
    sz=10, color=SEC, align=PP_ALIGN.CENTER)

snum(s4, 4)

# ============================================================
# SAVE
# ============================================================
output = "/projects/sandbox/Activity_1_Research_Methodology.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
