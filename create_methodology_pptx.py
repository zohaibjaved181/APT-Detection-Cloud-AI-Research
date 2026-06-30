"""
Research Methodology & Data Collection Presentation
Topic: Detection of APT in Cloud Networks Using AI
Professional design with dark theme
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Color palette - professional dark theme
BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
SURFACE = RGBColor(0x33, 0x41, 0x55)
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
PRIMARY_LT = RGBColor(0x60, 0xA5, 0xFA)
TEAL = RGBColor(0x08, 0x91, 0xB2)
TEAL_LT = RGBColor(0x22, 0xD3, 0xEE)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LT = RGBColor(0xA7, 0x8B, 0xFA)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_LT = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SECONDARY = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)



def set_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

def transparency(shape, val):
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    sf = spPr.find('.//' + qn('a:solidFill'))
    if sf is not None:
        clr = sf.find(qn('a:srgbClr'))
        if clr is not None:
            a = etree.SubElement(clr, qn('a:alpha'))
            a.set('val', str(int((1 - val) * 100000)))

def orb(slide, l, t, w, h, color, alpha=0.85):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    transparency(s, alpha)

def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=14, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tb

def decor(slide):
    orb(slide, Inches(10.5), Inches(-0.8), Inches(3), Inches(3), PRIMARY)
    orb(slide, Inches(-0.8), Inches(5.5), Inches(2.2), Inches(2.2), TEAL, 0.9)

def snum(slide, n, total=7):
    txt(slide, Inches(11.8), Inches(6.95), Inches(1.2), Inches(0.35),
        f"{n:02d} / {total:02d}", sz=9, color=MUTED, align=PP_ALIGN.RIGHT)



# ============================================================
# SLIDE 1: Activity 1 - Topic & Methodology Need
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)
decor(s1)
orb(s1, Inches(5), Inches(2.5), Inches(4), Inches(4), PURPLE, 0.93)

txt(s1, Inches(2.5), Inches(0.7), Inches(8), Inches(0.35),
    "ACTIVITY 1: RESEARCH METHODOLOGY SELECTION",
    sz=9, bold=True, color=PRIMARY_LT, align=PP_ALIGN.CENTER)

txt(s1, Inches(1.5), Inches(1.3), Inches(10.3), Inches(1.5),
    "Choosing the Right Research Methodology",
    sz=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s1, Inches(1.8), Inches(2.8), Inches(9.7), Inches(1.5),
    "For our study on detecting Advanced Persistent Threats in cloud networks "
    "through artificial intelligence, we need a methodology that allows us to "
    "measure, test, and validate system performance in a structured way.",
    sz=13, color=SECONDARY, align=PP_ALIGN.CENTER)

# Left card - Topic
rect(s1, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.7), CARD, PRIMARY)
txt(s1, Inches(0.8), Inches(4.5), Inches(5.4), Inches(0.4),
    "Our Research Topic", sz=13, bold=True, color=PRIMARY_LT)
txt(s1, Inches(0.8), Inches(5.0), Inches(5.4), Inches(2.0),
    "We are building an AI-based detection system that identifies "
    "sophisticated, multi-stage cyberattacks in cloud environments. "
    "The work involves training machine learning models on network traffic, "
    "testing against known attack scenarios, and measuring performance "
    "against existing tools. This is fundamentally an engineering and "
    "experimental problem requiring numerical evidence.",
    sz=10, color=MUTED)

# Right card - What we need
rect(s1, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.7), CARD, TEAL)
txt(s1, Inches(7.1), Inches(4.5), Inches(5.4), Inches(0.4),
    "What We Need From a Methodology", sz=13, bold=True, color=TEAL_LT)
needs = [
    "Run controlled experiments with measurable outcomes",
    "Statistical tools to compare AI models objectively",
    "Clear, reproducible results others can verify",
    "Structured approach to testing detection hypotheses",
]
for i, n in enumerate(needs):
    txt(s1, Inches(7.3), Inches(5.1 + i*0.45), Inches(5.2), Inches(0.4),
        f"  {n}", sz=10, color=MUTED)

snum(s1, 1)



# ============================================================
# SLIDE 2: Methodology Comparison
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
decor(s2)

txt(s2, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
    "COMPARING APPROACHES", sz=9, bold=True, color=PURPLE_LT)
txt(s2, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "Qualitative vs. Quantitative vs. Mixed Methods",
    sz=26, bold=True, color=WHITE)
txt(s2, Inches(0.6), Inches(1.5), Inches(10), Inches(0.4),
    "Understanding what each methodology offers and where it falls short for our research problem.",
    sz=11, color=MUTED)

# Three columns
methods = [
    ("Qualitative", PURPLE_LT, PURPLE,
     "Focuses on understanding experiences and meanings through interviews and thematic analysis.",
     "Exploring how security analysts perceive and respond to threats.",
     "We need hard numbers on detection rates, not subjective interpretation."),
    ("Quantitative", PRIMARY_LT, PRIMARY,
     "Relies on numerical data, statistical analysis, and measurable variables. Tests hypotheses.",
     "Our entire research revolves around measuring accuracy, comparing algorithms statistically.",
     "Does not capture human factors. We accept this since our focus is the AI system."),
    ("Mixed Methods", TEAL_LT, TEAL,
     "Combines quantitative experiments with qualitative insights from participants.",
     "When you also want to explore usability of detection dashboards.",
     "Adds complexity without clear benefit. Our question is about model performance."),
]

for i, (name, color, border, desc, works, why_not) in enumerate(methods):
    x = Inches(0.4 + i * 4.2)
    rect(s2, x, Inches(2.1), Inches(3.9), Inches(5.2), CARD, border)
    txt(s2, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.4),
        name, sz=14, bold=True, color=color)
    # Selected badge for Quantitative
    if i == 1:
        txt(s2, x + Inches(2.5), Inches(2.35), Inches(1.2), Inches(0.3),
            "SELECTED", sz=7, bold=True, color=GREEN_LT)
    txt(s2, x + Inches(0.2), Inches(2.8), Inches(3.5), Inches(1.0),
        desc, sz=9, color=MUTED)
    # When it works
    label = "WHY IT FITS" if i == 1 else "WHEN IT WORKS"
    txt(s2, x + Inches(0.2), Inches(3.8), Inches(3.5), Inches(0.3),
        label, sz=8, bold=True, color=GREEN_LT)
    txt(s2, x + Inches(0.2), Inches(4.1), Inches(3.5), Inches(0.9),
        works, sz=9, color=MUTED)
    # Limitation
    label2 = "LIMITATION" if i == 1 else "WHY NOT FOR US"
    lcolor = ORANGE if i == 1 else RED
    txt(s2, x + Inches(0.2), Inches(5.1), Inches(3.5), Inches(0.3),
        label2, sz=8, bold=True, color=lcolor)
    txt(s2, x + Inches(0.2), Inches(5.4), Inches(3.5), Inches(0.9),
        why_not, sz=9, color=MUTED)

snum(s2, 2)



# ============================================================
# SLIDE 3: Justification for Quantitative
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
decor(s3)

txt(s3, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
    "JUSTIFICATION", sz=9, bold=True, color=GREEN_LT)
txt(s3, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "Why Quantitative Methodology is the Best Fit",
    sz=26, bold=True, color=WHITE)
txt(s3, Inches(0.6), Inches(1.5), Inches(10), Inches(0.4),
    "The nature of our research problem naturally demands a quantitative approach.",
    sz=11, color=MUTED)

# Left - Core argument
rect(s3, Inches(0.4), Inches(2.1), Inches(6.2), Inches(4.3), CARD, PRIMARY)
txt(s3, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.4),
    "The Core Argument", sz=13, bold=True, color=PRIMARY_LT)
txt(s3, Inches(0.7), Inches(2.8), Inches(5.5), Inches(0.7),
    "At its heart, our research asks: Can an AI model reliably detect "
    "advanced persistent threats in cloud network traffic? Answering this requires:",
    sz=10, color=SECONDARY)

steps = [
    "Train models on large-scale datasets with thousands of flow records",
    "Measure detection using precision, recall, and F1-score",
    "Compare results against existing methods on same benchmarks",
    "Apply statistical tests to confirm improvements are not due to chance",
]
for i, step in enumerate(steps):
    txt(s3, Inches(0.9), Inches(3.6 + i*0.55), Inches(5.3), Inches(0.5),
        f"{i+1}.  {step}", sz=10, color=MUTED)

# Right - Alignment table
rect(s3, Inches(6.9), Inches(2.1), Inches(5.8), Inches(3.5), CARD, GREEN)
txt(s3, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.4),
    "Alignment with Research Goals", sz=13, bold=True, color=GREEN_LT)

goals = [
    ("Build detection framework", "Metrics evaluate if the system works"),
    ("Beat existing methods", "Statistical comparison gives objective evidence"),
    ("Ensure real-time viability", "Inference latency measured in milliseconds"),
    ("Validate across scenarios", "Cross-validation ensures generalizability"),
]
# Table header
rect(s3, Inches(7.0), Inches(2.8), Inches(5.5), Inches(0.4),
     RGBColor(0x15, 0x1D, 0x2F), None)
txt(s3, Inches(7.1), Inches(2.83), Inches(2.5), Inches(0.35),
    "Research Goal", sz=8, bold=True, color=PRIMARY_LT)
txt(s3, Inches(9.5), Inches(2.83), Inches(3.0), Inches(0.35),
    "How Quantitative Supports It", sz=8, bold=True, color=PRIMARY_LT)

for i, (goal, support) in enumerate(goals):
    y = Inches(3.25 + i*0.5)
    txt(s3, Inches(7.1), y, Inches(2.3), Inches(0.45),
        goal, sz=9, bold=True, color=WHITE)
    txt(s3, Inches(9.5), y, Inches(3.0), Inches(0.45),
        support, sz=9, color=MUTED)

# Bottom callout
rect(s3, Inches(6.9), Inches(5.8), Inches(5.8), Inches(1.1),
     RGBColor(0x0A, 0x1F, 0x20), GREEN)
txt(s3, Inches(7.2), Inches(5.9), Inches(5.2), Inches(0.9),
    "In plain terms: We chose quantitative methodology because our research "
    "produces numbers (accuracy, processing times, error rates), and we need "
    "statistical tools to make sense of those numbers and draw valid conclusions.",
    sz=10, color=SECONDARY)

snum(s3, 3)



# ============================================================
# SLIDE 4: Activity 1 Summary
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4)
decor(s4)

txt(s4, Inches(3.5), Inches(0.7), Inches(6), Inches(0.35),
    "ACTIVITY 1 SUMMARY", sz=9, bold=True, color=PRIMARY_LT, align=PP_ALIGN.CENTER)
txt(s4, Inches(2), Inches(1.3), Inches(9.3), Inches(0.8),
    "Methodology Decision: Quantitative",
    sz=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Visual indicator - 3 circles
for i, (label, color, sel) in enumerate([
    ("Quantitative", PRIMARY_LT, True),
    ("Qualitative", MUTED, False),
    ("Mixed", MUTED, False),
]):
    x = Inches(4.0 + i * 2.0)
    fill = RGBColor(0x15, 0x25, 0x45) if sel else RGBColor(0x1E, 0x29, 0x3B)
    bdr = PRIMARY if sel else None
    rect(s4, x, Inches(2.5), Inches(1.6), Inches(0.5), fill, bdr)
    mark = "  " + label if sel else "  " + label
    txt(s4, x, Inches(2.55), Inches(1.6), Inches(0.4),
        mark, sz=9, bold=sel, color=color, align=PP_ALIGN.CENTER)

# Key takeaway card
rect(s4, Inches(2.5), Inches(3.5), Inches(8.3), Inches(2.8), CARD, PRIMARY)
txt(s4, Inches(2.8), Inches(3.7), Inches(7.7), Inches(0.4),
    "Key Takeaway", sz=14, bold=True, color=PRIMARY_LT, align=PP_ALIGN.CENTER)
txt(s4, Inches(2.8), Inches(4.2), Inches(7.7), Inches(1.8),
    "A quantitative methodology gives us the scientific rigour needed to "
    "develop, test, and validate AI detection models. It lets us state precisely "
    "how well our system performs, back that up with statistical evidence, and "
    "present findings that the research community can independently verify "
    "and build upon.",
    sz=12, color=SECONDARY, align=PP_ALIGN.CENTER)

txt(s4, Inches(4), Inches(6.5), Inches(5.3), Inches(0.3),
    "Next: Activity 2 - Data Collection Methods",
    sz=10, color=MUTED, align=PP_ALIGN.CENTER)

snum(s4, 4)



# ============================================================
# SLIDE 5: Activity 2 - Data Collection Introduction
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5)
decor(s5)

txt(s5, Inches(0.6), Inches(0.4), Inches(5), Inches(0.3),
    "ACTIVITY 2: DATA COLLECTION METHODS", sz=9, bold=True, color=TEAL_LT)
txt(s5, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "How We Collect Our Data", sz=26, bold=True, color=WHITE)
txt(s5, Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
    "We have selected two complementary methods that together give us comprehensive coverage "
    "for training, testing, and validating our AI models.",
    sz=11, color=MUTED)

# Method 1
rect(s5, Inches(0.4), Inches(2.3), Inches(6.2), Inches(4.5), CARD, PRIMARY)
txt(s5, Inches(0.7), Inches(2.5), Inches(5.5), Inches(0.4),
    "Method 1: Secondary Data (Benchmark Datasets)", sz=13, bold=True, color=PRIMARY_LT)
txt(s5, Inches(0.7), Inches(3.0), Inches(5.5), Inches(1.2),
    "We use publicly available, peer-reviewed network traffic datasets that "
    "contain both normal activity and labelled APT attack samples. These are "
    "widely accepted in the cybersecurity research community and allow direct "
    "comparison with published results.",
    sz=10, color=MUTED)
datasets = ["CICIDS 2017/2018", "UNSW-NB15", "Unraveled APT Dataset", "DARPA TC"]
for i, ds in enumerate(datasets):
    txt(s5, Inches(0.9), Inches(4.3 + i*0.4), Inches(5.0), Inches(0.35),
        f"  {ds}", sz=10, color=SECONDARY)

# Method 2
rect(s5, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.5), CARD, TEAL)
txt(s5, Inches(7.2), Inches(2.5), Inches(5.2), Inches(0.4),
    "Method 2: Controlled Experiment (Cloud Simulation)", sz=13, bold=True, color=TEAL_LT)
txt(s5, Inches(7.2), Inches(3.0), Inches(5.2), Inches(1.2),
    "We set up a cloud testbed and simulate realistic APT attack campaigns "
    "using the MITRE ATT&CK framework. This generates fresh, cloud-native data "
    "that reflects modern threat behaviors not captured in older datasets.",
    sz=10, color=MUTED)
tools = ["AWS/Azure Cloud Testbed", "MITRE ATT&CK Framework", "Atomic Red Team Tools", "Custom Attack Scenarios"]
for i, t in enumerate(tools):
    txt(s5, Inches(7.4), Inches(4.3 + i*0.4), Inches(5.0), Inches(0.35),
        f"  {t}", sz=10, color=SECONDARY)

# Why two methods callout
rect(s5, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.0), CARD, None)

snum(s5, 5)



# ============================================================
# SLIDE 6: Justification for Each Method
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6)
decor(s6)

txt(s6, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
    "DEEP DIVE", sz=9, bold=True, color=PRIMARY_LT)
txt(s6, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "Justifying Each Collection Method", sz=26, bold=True, color=WHITE)
txt(s6, Inches(0.6), Inches(1.5), Inches(10), Inches(0.4),
    "Each method addresses a different gap. Together, they produce a complete picture.",
    sz=11, color=MUTED)

# Method 1 justification
rect(s6, Inches(0.4), Inches(2.1), Inches(6.2), Inches(4.6), CARD, PRIMARY)
txt(s6, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.4),
    "Secondary Data - Why It Works", sz=13, bold=True, color=PRIMARY_LT)

m1_reasons = [
    ("Widely validated", "used by hundreds of studies, enabling direct comparison"),
    ("Pre-labelled", "attack types already annotated, reducing manual effort and bias"),
    ("Large scale", "millions of network flows for deep learning training"),
    ("Reproducible", "other researchers can replicate experiments using same data"),
]
for i, (title, desc) in enumerate(m1_reasons):
    y = Inches(2.9 + i*0.65)
    txt(s6, Inches(0.9), y, Inches(5.3), Inches(0.3),
        f"  {title}", sz=10, bold=True, color=WHITE)
    txt(s6, Inches(0.9), y + Inches(0.25), Inches(5.3), Inches(0.3),
        f"    {desc}", sz=9, color=MUTED)

rect(s6, Inches(0.5), Inches(5.6), Inches(5.9), Inches(0.9),
     RGBColor(0x0A, 0x1A, 0x2F), PRIMARY)
txt(s6, Inches(0.7), Inches(5.7), Inches(5.5), Inches(0.7),
    "These datasets serve as our training foundation and allow benchmarking "
    "against state-of-the-art using same evaluation conditions.",
    sz=9, color=SECONDARY)

# Method 2 justification
rect(s6, Inches(6.9), Inches(2.1), Inches(5.8), Inches(4.6), CARD, TEAL)
txt(s6, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.4),
    "Controlled Experiment - Why It Works", sz=13, bold=True, color=TEAL_LT)

m2_reasons = [
    ("Cloud-native", "captures API abuse, container escape, lateral movement"),
    ("Current threats", "simulates 2024-2025 APT tactics older datasets miss"),
    ("Controlled variables", "we decide exactly which attacks, when, and intensity"),
    ("Fills research gap", "addresses documented lack of cloud-specific APT data"),
]
for i, (title, desc) in enumerate(m2_reasons):
    y = Inches(2.9 + i*0.65)
    txt(s6, Inches(7.4), y, Inches(5.0), Inches(0.3),
        f"  {title}", sz=10, bold=True, color=WHITE)
    txt(s6, Inches(7.4), y + Inches(0.25), Inches(5.0), Inches(0.3),
        f"    {desc}", sz=9, color=MUTED)

rect(s6, Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.9),
     RGBColor(0x0A, 0x1F, 0x20), TEAL)
txt(s6, Inches(7.2), Inches(5.7), Inches(5.1), Inches(0.7),
    "Simulated environment validates our model works on fresh, realistic "
    "cloud attack scenarios representing current threats.",
    sz=9, color=SECONDARY)

snum(s6, 6)



# ============================================================
# SLIDE 7: Final Summary
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7)
decor(s7)

txt(s7, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
    "PUTTING IT ALL TOGETHER", sz=9, bold=True, color=GREEN_LT)
txt(s7, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "Complete Research Design Summary", sz=26, bold=True, color=WHITE)
txt(s7, Inches(0.6), Inches(1.5), Inches(10), Inches(0.4),
    "A coherent design where methodology and data collection work together to answer our questions.",
    sz=11, color=MUTED)

# Left - Design table
rect(s7, Inches(0.4), Inches(2.1), Inches(6.2), Inches(3.5), CARD, PRIMARY)
txt(s7, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.4),
    "Design at a Glance", sz=13, bold=True, color=PRIMARY_LT)

design = [
    ("Research Methodology", "Quantitative"),
    ("Data Method 1", "Secondary Data (Benchmarks)"),
    ("Data Method 2", "Controlled Experiment (Cloud Sim)"),
    ("Analysis Approach", "Statistical + ML Metrics"),
    ("Validation", "Cross-validation + Hypothesis Tests"),
]
for i, (comp, choice) in enumerate(design):
    y = Inches(2.85 + i*0.5)
    bg = RGBColor(0x15, 0x1D, 0x2F) if i % 2 == 0 else RGBColor(0x12, 0x1A, 0x2A)
    rect(s7, Inches(0.5), y, Inches(5.9), Inches(0.45), bg, None)
    txt(s7, Inches(0.7), y + Inches(0.05), Inches(2.8), Inches(0.35),
        comp, sz=10, bold=True, color=WHITE)
    txt(s7, Inches(3.5), y + Inches(0.05), Inches(2.8), Inches(0.35),
        choice, sz=10, color=TEAL_LT)

# Right - Why it works
rect(s7, Inches(6.9), Inches(2.1), Inches(5.8), Inches(3.5), CARD, GREEN)
txt(s7, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.4),
    "Why This Design Works", sz=13, bold=True, color=GREEN_LT)

works = [
    "Quantitative methodology matches our numerical, performance-driven questions",
    "Benchmark datasets allow fair comparison with existing published methods",
    "Simulated environment validates real-world cloud applicability",
    "Together, they produce results that are academically rigorous and practical",
]
for i, w in enumerate(works):
    txt(s7, Inches(7.4), Inches(2.9 + i*0.6), Inches(5.0), Inches(0.5),
        f"  {w}", sz=10, color=MUTED)

# Bottom callout
rect(s7, Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.2),
     RGBColor(0x0A, 0x1F, 0x20), GREEN)
txt(s7, Inches(0.7), Inches(6.05), Inches(12.0), Inches(1.0),
    "Bottom line: Every piece of our research design connects logically. We ask a quantitative "
    "question, collect quantitative data through two complementary channels, and analyse it with "
    "statistical tools that give us clear, defensible answers.",
    sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

snum(s7, 7)

# ============================================================
# SAVE
# ============================================================
output = "/projects/sandbox/Research_Methodology_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
