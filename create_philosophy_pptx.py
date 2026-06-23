"""
Generate a modern PPTX presentation on Research Philosophy & Approach
for: Detection of APT in Cloud Networks Using AI
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Color palette
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
PRIMARY = RGBColor(0x63, 0x66, 0xF1)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_PRIMARY = RGBColor(0x81, 0x8C, 0xF8)
LIGHT_ACCENT = RGBColor(0x22, 0xD3, 0xEE)
LIGHT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
LIGHT_GREEN = RGBColor(0x4A, 0xDE, 0x80)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_transparency(shape, transparency):
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        clr = solidFill.find(qn('a:srgbClr'))
        if clr is not None:
            alpha = etree.SubElement(clr, qn('a:alpha'))
            alpha.set('val', str(int((1 - transparency) * 100000)))


def add_orb(slide, left, top, width, height, color, transparency=0.85):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    set_shape_transparency(shape, transparency)


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines):
    """lines: list of (text, size, bold, color)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
    return txBox


def decorate(slide):
    add_orb(slide, Inches(10), Inches(-0.5), Inches(3.5), Inches(3.5), PRIMARY)
    add_orb(slide, Inches(-1), Inches(5), Inches(2.5), Inches(2.5), ACCENT, 0.88)


def slide_num(slide, n):
    add_text(slide, Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4),
             f"{n:02d} / 05", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: Research Topic
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s1, DARK_BG)
decorate(s1)
add_orb(s1, Inches(4), Inches(2), Inches(5), Inches(5), PURPLE, 0.92)

add_text(s1, Inches(2), Inches(0.8), Inches(9), Inches(0.4),
         "RESEARCH METHODOLOGY  |  PHILOSOPHY & APPROACH",
         size=10, bold=True, color=LIGHT_PRIMARY, align=PP_ALIGN.CENTER)

add_text(s1, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.0),
         "Detection of Advanced Persistent Threats\nin Cloud Networks Using Artificial Intelligence",
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s1, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.8),
         "This research investigates the application of AI and machine learning techniques to detect "
         "sophisticated, multi-stage APTs targeting cloud infrastructure. The study aims to develop a "
         "novel, explainable detection framework addressing gaps in real-time processing, adversarial "
         "robustness, and cross-stage threat correlation.",
         size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Highlight box
add_rect(s1, Inches(2.0), Inches(5.3), Inches(9.3), Inches(1.2),
         RGBColor(0x0A, 0x1F, 0x2F), ACCENT)
add_text(s1, Inches(2.3), Inches(5.45), Inches(8.7), Inches(1.0),
         "Research Focus: Empirical evaluation of AI models on cloud network data, measuring "
         "detection accuracy, false positive rates, inference latency, and explainability of "
         "threat classifications across multi-cloud environments.",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

slide_num(s1, 1)

# ============================================================
# SLIDE 2: Research Philosophies
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s2, DARK_BG)
decorate(s2)

add_text(s2, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
         "EPISTEMOLOGY", size=9, bold=True, color=LIGHT_PURPLE)
add_text(s2, Inches(0.6), Inches(0.8), Inches(10), Inches(0.7),
         "Research Philosophies", size=28, bold=True, color=WHITE)
add_text(s2, Inches(0.6), Inches(1.4), Inches(10), Inches(0.5),
         "Four major philosophies and their applicability to AI-based APT detection research",
         size=12, color=MUTED)

philosophies = [
    ("Positivism", PRIMARY, LIGHT_PRIMARY,
     "Knowledge through observable, measurable facts. Emphasizes objectivity and hypothesis testing.",
     "Enables quantitative evaluation (accuracy, precision, recall). Reproducible experiments.",
     "May overlook contextual/adaptive nature of APT actors."),
    ("Interpretivism", ACCENT, LIGHT_ACCENT,
     "Understanding through subjective experiences and meanings. Focuses on context and interpretation.",
     "Useful for understanding SOC analyst decision-making and alert interpretation.",
     "Not ideal for validating AI model performance. Lacks generalizability."),
    ("Pragmatism", PURPLE, LIGHT_PURPLE,
     "Knowledge judged by practical consequences. Supports mixed methods approach.",
     "Combines quantitative metrics with qualitative expert validation. Flexible.",
     "May lack philosophical rigour. Risk of paradigm inconsistency."),
    ("Realism", GREEN, LIGHT_GREEN,
     "Reality exists independently. Accepts observable data and underlying causal mechanisms.",
     "Acknowledges complex, layered nature of APTs. Explores causal mechanisms.",
     "Harder to operationalize in purely computational experiments."),
]

for i, (name, border, text_color, desc, strength, limitation) in enumerate(philosophies):
    x = Inches(0.4 + i * 3.2)
    add_rect(s2, x, Inches(2.0), Inches(3.0), Inches(5.2),
             RGBColor(0x1E, 0x29, 0x3B), border)
    add_text(s2, x + Inches(0.2), Inches(2.2), Inches(2.6), Inches(0.4),
             name, size=13, bold=True, color=text_color)
    add_text(s2, x + Inches(0.2), Inches(2.7), Inches(2.6), Inches(1.2),
             desc, size=9, color=MUTED)
    add_text(s2, x + Inches(0.2), Inches(3.8), Inches(2.6), Inches(0.3),
             "STRENGTHS", size=8, bold=True, color=LIGHT_GREEN)
    add_text(s2, x + Inches(0.2), Inches(4.1), Inches(2.6), Inches(1.0),
             strength, size=9, color=MUTED)
    add_text(s2, x + Inches(0.2), Inches(5.1), Inches(2.6), Inches(0.3),
             "LIMITATIONS", size=8, bold=True, color=RED)
    add_text(s2, x + Inches(0.2), Inches(5.4), Inches(2.6), Inches(1.0),
             limitation, size=9, color=MUTED)

slide_num(s2, 2)

# ============================================================
# SLIDE 3: Selected Philosophy
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s3, DARK_BG)
decorate(s3)

add_text(s3, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         "SELECTED PHILOSOPHY", size=9, bold=True, color=LIGHT_GREEN)
add_text(s3, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
         "Positivism - The Most Suitable Philosophy", size=28, bold=True, color=WHITE)
add_text(s3, Inches(0.6), Inches(1.5), Inches(10), Inches(0.4),
         "Positivism is selected as the guiding research philosophy for this AI-based APT detection study",
         size=12, color=MUTED)

# Left - Why Positivism
add_rect(s3, Inches(0.4), Inches(2.1), Inches(6.2), Inches(4.8),
         RGBColor(0x1E, 0x29, 0x3B), PRIMARY)
add_text(s3, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.4),
         "Why Positivism?", size=14, bold=True, color=LIGHT_PRIMARY)

reasons = [
    "Research relies on quantitative, measurable outcomes - accuracy, precision, recall, F1-score",
    "Requires objective, reproducible experiments using standardized datasets",
    "Involves hypothesis testing - 'AI model X detects APTs better than baseline'",
    "Supports generalizability across different cloud environments",
    "Aligns with scientific method dominant in CS and cybersecurity research",
]
for i, reason in enumerate(reasons):
    add_text(s3, Inches(0.9), Inches(2.9 + i * 0.7), Inches(5.5), Inches(0.6),
             f"  {reason}", size=10, color=MUTED)

# Right - Alignment
add_rect(s3, Inches(6.9), Inches(2.1), Inches(5.8), Inches(3.2),
         RGBColor(0x1E, 0x29, 0x3B), ACCENT)
add_text(s3, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.4),
         "Alignment with Research Problem", size=14, bold=True, color=LIGHT_ACCENT)

alignments = [
    "Empirical validation of AI model performance on network traffic data",
    "Statistical comparison of multiple detection algorithms",
    "Controlled experiments with train/test splits and cross-validation",
    "Value-free observation - system performance is objective",
    "Identifying causal relationships between features and detection",
]
for i, a in enumerate(alignments):
    add_text(s3, Inches(7.4), Inches(2.8 + i * 0.5), Inches(5.0), Inches(0.45),
             f"  {a}", size=10, color=MUTED)

# Bottom insight box
add_rect(s3, Inches(6.9), Inches(5.6), Inches(5.8), Inches(1.3),
         RGBColor(0x0A, 0x1F, 0x2F), ACCENT)
add_text(s3, Inches(7.2), Inches(5.75), Inches(5.2), Inches(1.0),
         "Key Insight: Since APT detection is fundamentally a classification/anomaly detection "
         "problem evaluated through numerical metrics, positivism provides the most rigorous "
         "framework for validating our AI approach.",
         size=10, color=MUTED)

slide_num(s3, 3)

# ============================================================
# SLIDE 4: Research Approach
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s4, DARK_BG)
decorate(s4)

add_text(s4, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         "RESEARCH APPROACH", size=9, bold=True, color=LIGHT_ACCENT)
add_text(s4, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
         "Deductive Approach - Theory to Evidence", size=28, bold=True, color=WHITE)
add_text(s4, Inches(0.6), Inches(1.5), Inches(10), Inches(0.4),
         "Moving from established theories and hypotheses to empirical testing and validation",
         size=12, color=MUTED)

# Left - Deductive vs Inductive comparison
add_rect(s4, Inches(0.4), Inches(2.1), Inches(6.2), Inches(1.4),
         RGBColor(0x1E, 0x29, 0x3B), PURPLE)
# Deductive box
add_rect(s4, Inches(0.6), Inches(2.3), Inches(2.8), Inches(1.0),
         RGBColor(0x15, 0x1D, 0x3F), PRIMARY)
add_text(s4, Inches(0.8), Inches(2.35), Inches(2.4), Inches(0.3),
         "DEDUCTIVE  [Selected]", size=9, bold=True, color=LIGHT_PRIMARY)
add_text(s4, Inches(0.8), Inches(2.7), Inches(2.4), Inches(0.5),
         "Theory > Hypothesis > Data > Testing > Confirm/Reject", size=8, color=MUTED)
# Inductive box
add_rect(s4, Inches(3.6), Inches(2.3), Inches(2.8), Inches(1.0),
         RGBColor(0x1F, 0x15, 0x15), RED)
add_text(s4, Inches(3.8), Inches(2.35), Inches(2.4), Inches(0.3),
         "INDUCTIVE  [Not selected]", size=9, bold=True, color=RED)
add_text(s4, Inches(3.8), Inches(2.7), Inches(2.4), Inches(0.5),
         "Observation > Pattern > Hypothesis > Theory (exploratory)", size=8, color=MUTED)

# Left - Deductive Process Flow
add_rect(s4, Inches(0.4), Inches(3.7), Inches(6.2), Inches(3.5),
         RGBColor(0x1E, 0x29, 0x3B), ACCENT)
add_text(s4, Inches(0.7), Inches(3.9), Inches(5.5), Inches(0.4),
         "Deductive Process for This Study", size=13, bold=True, color=LIGHT_ACCENT)

steps = [
    ("1. Theory:", "AI/DL can model complex, non-linear attack patterns", PRIMARY),
    ("2. Hypothesis:", "Proposed hybrid model achieves >95% detection rate", ACCENT),
    ("3. Experiment:", "Train/test on cloud APT datasets with baselines", PURPLE),
    ("4. Validate:", "Statistical tests confirm or reject hypothesis", GREEN),
]
for i, (label, desc, color) in enumerate(steps):
    y = Inches(4.4 + i * 0.7)
    add_rect(s4, Inches(0.7), y, Inches(5.6), Inches(0.55), RGBColor(0x12, 0x1A, 0x2A), color)
    add_text(s4, Inches(0.9), y + Inches(0.05), Inches(1.5), Inches(0.4),
             label, size=10, bold=True, color=color)
    add_text(s4, Inches(2.4), y + Inches(0.05), Inches(3.7), Inches(0.4),
             desc, size=10, color=MUTED)

# Right - Why Deductive
add_rect(s4, Inches(6.9), Inches(2.1), Inches(5.8), Inches(5.1),
         RGBColor(0x1E, 0x29, 0x3B), GREEN)
add_text(s4, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.4),
         "Why Deductive is Suitable", size=14, bold=True, color=LIGHT_GREEN)

deductive_reasons = [
    ("1", "Existing theoretical foundation in ML/DL pattern recognition provides clear starting hypotheses"),
    ("2", "Testable predictions - specific, falsifiable hypotheses about model performance metrics"),
    ("3", "Structured methodology - systematic comparison of proposed vs. existing approaches"),
    ("4", "Generalizability - results applicable beyond specific test dataset to broader environments"),
    ("5", "Alignment with positivism - deductive naturally complements positivist philosophy"),
    ("6", "Benchmarking standard - consistent with how ML/AI research is conducted in top venues"),
]
for i, (num, reason) in enumerate(deductive_reasons):
    y = Inches(2.85 + i * 0.7)
    add_text(s4, Inches(7.2), y, Inches(0.4), Inches(0.35),
             num, size=11, bold=True, color=LIGHT_GREEN)
    add_text(s4, Inches(7.6), y, Inches(4.9), Inches(0.6),
             reason, size=10, color=MUTED)

slide_num(s4, 4)

# ============================================================
# SLIDE 5: Final Justification
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s5, DARK_BG)
decorate(s5)

add_text(s5, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         "FINAL JUSTIFICATION", size=9, bold=True, color=LIGHT_GREEN)
add_text(s5, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
         "How Philosophy & Approach Support the Research", size=28, bold=True, color=WHITE)
add_text(s5, Inches(0.6), Inches(1.5), Inches(10), Inches(0.4),
         "Coherent alignment between Positivism + Deductive approach and objectives, data collection, and analysis",
         size=12, color=MUTED)

# Left - Research Objectives
add_rect(s5, Inches(0.4), Inches(2.1), Inches(6.2), Inches(3.5),
         RGBColor(0x1E, 0x29, 0x3B), PRIMARY)
add_text(s5, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.4),
         "Supporting Research Objectives", size=14, bold=True, color=LIGHT_PRIMARY)

objectives = [
    ("Obj 1: Develop AI detection framework",
     "Positivism enables objective measurement of framework performance"),
    ("Obj 2: Compare with existing methods",
     "Deductive approach structures hypothesis-driven comparison"),
    ("Obj 3: Validate in cloud environments",
     "Positivism demands empirical testing in realistic scenarios"),
    ("Obj 4: Ensure explainability",
     "Measurable XAI metrics (fidelity, stability) align with positivist evaluation"),
]
for i, (obj, support) in enumerate(objectives):
    y = Inches(2.8 + i * 0.8)
    add_text(s5, Inches(0.9), y, Inches(5.5), Inches(0.35),
             obj, size=10, bold=True, color=WHITE)
    add_text(s5, Inches(0.9), y + Inches(0.3), Inches(5.5), Inches(0.35),
             f"  {support}", size=9, color=LIGHT_GREEN)

# Right top - Data Collection
add_rect(s5, Inches(6.9), Inches(2.1), Inches(5.8), Inches(2.0),
         RGBColor(0x1E, 0x29, 0x3B), ACCENT)
add_text(s5, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.4),
         "Supporting Data Collection", size=14, bold=True, color=LIGHT_ACCENT)
dc_items = [
    "Quantitative network traffic datasets (CICIDS, Unraveled, DARPA)",
    "Structured numerical features from cloud logs & flows",
    "Controlled experimental conditions with train/val/test splits",
    "Reproducible simulations using MITRE ATT&CK scenarios",
]
for i, item in enumerate(dc_items):
    add_text(s5, Inches(7.4), Inches(2.8 + i * 0.35), Inches(5.0), Inches(0.3),
             f"  {item}", size=9, color=MUTED)

# Right bottom - Data Analysis
add_rect(s5, Inches(6.9), Inches(4.3), Inches(5.8), Inches(2.0),
         RGBColor(0x1E, 0x29, 0x3B), PURPLE)
add_text(s5, Inches(7.2), Inches(4.5), Inches(5.2), Inches(0.4),
         "Supporting Data Analysis", size=14, bold=True, color=LIGHT_PURPLE)
da_items = [
    "Statistical hypothesis testing (t-tests, ANOVA, Wilcoxon)",
    "Performance metrics - Accuracy, Precision, Recall, F1, AUC-ROC",
    "Confusion matrices & ROC curves for visual validation",
    "Cross-validation (k-fold) ensuring robustness & generalizability",
]
for i, item in enumerate(da_items):
    add_text(s5, Inches(7.4), Inches(5.0 + i * 0.35), Inches(5.0), Inches(0.3),
             f"  {item}", size=9, color=MUTED)

# Bottom conclusion
add_rect(s5, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.85),
         RGBColor(0x0A, 0x1F, 0x2F), ACCENT)
add_text(s5, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.7),
         "Conclusion: The combination of Positivism (philosophy) and Deductive Approach provides a rigorous, "
         "objective, and systematic framework perfectly suited to developing, testing, and validating an "
         "AI-based APT detection system through empirical experimentation and statistical analysis.",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_num(s5, 5)

# ============================================================
# Save
# ============================================================
output_path = "/projects/sandbox/Research_Philosophy_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
