"""
Generate a modern PPTX presentation on:
Detection of Advanced Persistent Threats in Cloud Networks Using AI
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
import math

# Color palette
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARKER_BG = RGBColor(0x02, 0x06, 0x17)
PRIMARY = RGBColor(0x63, 0x66, 0xF1)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_PRIMARY = RGBColor(0x81, 0x8C, 0xF8)
LIGHT_ACCENT = RGBColor(0x22, 0xD3, 0xEE)
LIGHT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_transparency(shape, transparency):
    """Set transparency on a shape's solid fill using XML manipulation."""
    from lxml import etree
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        clr = solidFill.find(qn('a:srgbClr'))
        if clr is not None:
            alpha = etree.SubElement(clr, qn('a:alpha'))
            alpha.set('val', str(int((1 - transparency) * 100000)))

def add_gradient_rect(slide, left, top, width, height, color1, color2, transparency=0.7):
    """Add a gradient rectangle as decorative element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    set_shape_transparency(shape, transparency)

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=14):
    """Add text box with multiple formatted lines.
    lines: list of tuples (text, font_size, bold, color)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text, size, bold, color = line_data
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, transparency=0.0):
    """Add rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if transparency > 0:
        set_shape_transparency(shape, transparency)
    return shape


def add_card_with_text(slide, left, top, width, height, title, body_lines,
                       title_color=LIGHT_PRIMARY, card_color=RGBColor(0x1E, 0x29, 0x3B)):
    """Add a card with title and bullet points."""
    shape = add_rounded_rect(slide, left, top, width, height, card_color, PRIMARY)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4),
                 Inches(0.4), title, font_size=13, bold=True, color=title_color)
    # Body
    body_top = top + Inches(0.55)
    lines = [(line, 11, False, MUTED) for line in body_lines]
    add_multi_text(slide, left + Inches(0.2), body_top, width - Inches(0.4),
                   height - Inches(0.7), lines, default_size=11)

def add_slide_number(slide, num):
    """Add slide number to bottom right."""
    add_text_box(slide, Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4),
                 f"{num:02d} / 10", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)

def add_tag(slide, left, top, text, color=LIGHT_PRIMARY):
    """Add a small tag/badge."""
    add_text_box(slide, left, top, Inches(2), Inches(0.3), text.upper(),
                 font_size=9, bold=True, color=color)

def decorate_slide(slide):
    """Add decorative gradient orbs to a slide."""
    add_gradient_rect(slide, Inches(10), Inches(-0.5), Inches(3.5), Inches(3.5), PRIMARY, ACCENT, 0.85)
    add_gradient_rect(slide, Inches(-1), Inches(5), Inches(2.5), Inches(2.5), ACCENT, PURPLE, 0.88)


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1, DARK_BG)
decorate_slide(slide1)
add_gradient_rect(slide1, Inches(5), Inches(2), Inches(4), Inches(4), PURPLE, PRIMARY, 0.9)

# Tags
add_tag(slide1, Inches(3.5), Inches(1.0), "RESEARCH PRESENTATION", LIGHT_PRIMARY)
add_tag(slide1, Inches(6.5), Inches(1.0), "2024-2025", LIGHT_ACCENT)

# Title
add_text_box(slide1, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.5),
             "Detection of Advanced Persistent Threats\nin Cloud Networks Using Artificial Intelligence",
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(2), Inches(3.8), Inches(9.3), Inches(1.0),
             "A comprehensive research study exploring AI-driven methodologies for identifying\nand mitigating sophisticated cyber threats in cloud computing environments",
             font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# Stats boxes
stats = [("71%", "Organizations Hit by APTs"), ("197", "Avg. Days to Detect"), ("$4.9M", "Avg. Breach Cost")]
for i, (val, label) in enumerate(stats):
    x = Inches(2.5 + i * 3.0)
    add_rounded_rect(slide1, x, Inches(5.2), Inches(2.5), Inches(1.5),
                     RGBColor(0x1E, 0x29, 0x3B), PRIMARY)
    add_text_box(slide1, x, Inches(5.35), Inches(2.5), Inches(0.7),
                 val, font_size=28, bold=True, color=LIGHT_ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, x, Inches(6.0), Inches(2.5), Inches(0.5),
                 label, font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide1, 1)


# ============================================================
# SLIDE 2: Introduction to APTs
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, DARK_BG)
decorate_slide(slide2)

add_tag(slide2, Inches(0.6), Inches(0.4), "THREAT LANDSCAPE", RED)
add_text_box(slide2, Inches(0.6), Inches(0.8), Inches(8), Inches(0.8),
             "What Are Advanced Persistent Threats?",
             font_size=28, bold=True, color=WHITE)

# Description
add_text_box(slide2, Inches(0.6), Inches(1.6), Inches(6), Inches(1.2),
             "APTs are prolonged, targeted cyberattacks where adversaries gain unauthorized "
             "network access and remain undetected for extended periods. They target high-value "
             "assets through sophisticated multi-stage campaigns.",
             font_size=12, color=MUTED)

# Characteristics tags
chars = ["Stealthy", "Persistent", "Targeted", "Multi-Stage", "Adaptive"]
char_colors = [RED, PRIMARY, ACCENT, PURPLE, GREEN]
for i, (c, clr) in enumerate(zip(chars, char_colors)):
    add_text_box(slide2, Inches(0.6 + i * 1.6), Inches(2.8), Inches(1.5), Inches(0.3),
                 f"  {c}  ", font_size=9, bold=True, color=clr)

# Left card - Notable APT Groups
add_card_with_text(slide2, Inches(0.6), Inches(3.3), Inches(5.8), Inches(1.6),
                   "Notable APT Groups",
                   ["APT28 (Fancy Bear) - Russian state-sponsored",
                    "APT29 (Cozy Bear) - SolarWinds attack",
                    "Lazarus Group - North Korean operations",
                    "APT41 - Chinese dual espionage/cybercrime"],
                   title_color=LIGHT_PRIMARY)

# Right card - Why Traditional Defenses Fail
add_card_with_text(slide2, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.3),
                   "Why Traditional Defenses Fail",
                   ["X  Signature-based IDS cannot detect zero-day exploits",
                    "X  Firewall rules bypassed via encrypted channels",
                    "X  Low-and-slow tactics evade threshold alerts",
                    "X  Lateral movement mimics legitimate traffic",
                    "X  Polymorphic malware evades static analysis"],
                   title_color=RED)

# Stat callout
add_rounded_rect(slide2, Inches(0.6), Inches(5.3), Inches(12), Inches(1.2),
                 RGBColor(0x2D, 0x1B, 0x1B), RED)
add_text_box(slide2, Inches(1.0), Inches(5.5), Inches(3), Inches(0.7),
             "+150% Increase", font_size=22, bold=True, color=RED)
add_text_box(slide2, Inches(1.0), Inches(6.0), Inches(10), Inches(0.4),
             "In APT attacks on cloud infrastructure (2022-2024)", font_size=11, color=MUTED)

add_slide_number(slide2, 2)


# ============================================================
# SLIDE 3: APT Attack Lifecycle
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, DARK_BG)
decorate_slide(slide3)

add_tag(slide3, Inches(0.6), Inches(0.4), "KILL CHAIN MODEL", LIGHT_PURPLE)
add_text_box(slide3, Inches(0.6), Inches(0.8), Inches(8), Inches(0.8),
             "APT Attack Lifecycle", font_size=28, bold=True, color=WHITE)
add_text_box(slide3, Inches(0.6), Inches(1.5), Inches(9), Inches(0.5),
             "Based on Lockheed Martin Cyber Kill Chain and MITRE ATT&CK Framework",
             font_size=12, color=MUTED)

# Kill chain flow - 7 stages as colored boxes
stages = [
    ("Reconnaissance", RGBColor(0x3D, 0x1F, 0x1F), RED),
    ("Weaponization", RGBColor(0x3D, 0x2A, 0x1A), ORANGE),
    ("Delivery", RGBColor(0x3D, 0x3A, 0x1A), YELLOW),
    ("Exploitation", RGBColor(0x1A, 0x3D, 0x1F), GREEN),
    ("Installation", RGBColor(0x1A, 0x2E, 0x3D), LIGHT_ACCENT),
    ("C2", RGBColor(0x1F, 0x1F, 0x3D), LIGHT_PRIMARY),
    ("Exfiltration", RGBColor(0x2A, 0x1F, 0x3D), LIGHT_PURPLE),
]

for i, (name, bg, text_color) in enumerate(stages):
    x = Inches(0.4 + i * 1.82)
    add_rounded_rect(slide3, x, Inches(2.2), Inches(1.65), Inches(0.7), bg, text_color)
    add_text_box(slide3, x, Inches(2.3), Inches(1.65), Inches(0.6),
                 name, font_size=10, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
    if i < 6:
        add_text_box(slide3, x + Inches(1.55), Inches(2.35), Inches(0.4), Inches(0.4),
                     ">", font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# Three info cards below
cards_data = [
    ("Initial Access", "Spear phishing, supply chain compromise,\nzero-day exploits targeting cloud APIs\nand misconfigured services", RED),
    ("Persistence & Lateral Movement", "Privilege escalation, credential harvesting,\nmoving across cloud workloads\nand containers undetected", PRIMARY),
    ("Command & Control + Exfiltration", "Encrypted C2 channels via legitimate\ncloud services (DNS tunneling, HTTPS\nbeaconing), slow data exfiltration", PURPLE),
]

for i, (title, body, color) in enumerate(cards_data):
    x = Inches(0.4 + i * 4.2)
    add_rounded_rect(slide3, x, Inches(3.3), Inches(3.9), Inches(3.2),
                     RGBColor(0x1E, 0x29, 0x3B), color)
    add_text_box(slide3, x + Inches(0.2), Inches(3.5), Inches(3.5), Inches(0.5),
                 title, font_size=13, bold=True, color=color)
    add_text_box(slide3, x + Inches(0.2), Inches(4.1), Inches(3.5), Inches(2.2),
                 body, font_size=11, color=MUTED)

add_slide_number(slide3, 3)


# ============================================================
# SLIDE 4: Cloud-Specific Threat Vectors
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, DARK_BG)
decorate_slide(slide4)

add_tag(slide4, Inches(0.6), Inches(0.4), "CLOUD SECURITY", LIGHT_ACCENT)
add_text_box(slide4, Inches(0.6), Inches(0.8), Inches(8), Inches(0.8),
             "Cloud-Specific Threat Vectors", font_size=28, bold=True, color=WHITE)
add_text_box(slide4, Inches(0.6), Inches(1.5), Inches(9), Inches(0.5),
             "Cloud environments introduce unique attack surfaces that APT actors actively exploit",
             font_size=12, color=MUTED)

# Left column - Attack Surfaces
surfaces = [
    ("API Exploitation", "Misconfigured REST APIs, exposed credentials in metadata services"),
    ("Container Escape", "Kubernetes vulnerabilities, Docker breakout to host system"),
    ("Identity Compromise", "IAM role abuse, OAuth token theft, federation attacks"),
    ("Serverless Injection", "Lambda/Function poisoning, event-driven trigger manipulation"),
]

add_rounded_rect(slide4, Inches(0.4), Inches(2.1), Inches(6.0), Inches(4.8),
                 RGBColor(0x1E, 0x29, 0x3B), ACCENT)
add_text_box(slide4, Inches(0.7), Inches(2.3), Inches(5), Inches(0.4),
             "Cloud Attack Surfaces", font_size=14, bold=True, color=LIGHT_ACCENT)

for i, (title, desc) in enumerate(surfaces):
    y = Inches(2.9 + i * 1.1)
    add_text_box(slide4, Inches(0.9), y, Inches(5.2), Inches(0.35),
                 title, font_size=12, bold=True, color=WHITE)
    add_text_box(slide4, Inches(0.9), y + Inches(0.3), Inches(5.2), Inches(0.5),
                 desc, font_size=10, color=MUTED)

# Right column - Detection Challenges
challenges = [
    "Multi-tenancy obscures attack attribution",
    "Ephemeral workloads destroy forensic evidence",
    "Encrypted east-west traffic limits visibility",
    "Dynamic scaling creates baseline drift",
    "Shared responsibility model creates gaps",
]

add_rounded_rect(slide4, Inches(6.8), Inches(2.1), Inches(5.8), Inches(3.2),
                 RGBColor(0x1E, 0x29, 0x3B), PURPLE)
add_text_box(slide4, Inches(7.1), Inches(2.3), Inches(5), Inches(0.4),
             "Detection Challenges in Cloud", font_size=14, bold=True, color=LIGHT_PURPLE)

for i, ch in enumerate(challenges):
    y = Inches(2.8 + i * 0.5)
    bullet_colors = [RED, ORANGE, YELLOW, GREEN, LIGHT_PRIMARY]
    add_text_box(slide4, Inches(7.1), y, Inches(5.2), Inches(0.4),
                 f"  {ch}", font_size=11, color=MUTED)

# Stats box
add_rounded_rect(slide4, Inches(6.8), Inches(5.6), Inches(5.8), Inches(1.3),
                 RGBColor(0x1A, 0x2E, 0x3D), ACCENT)
add_text_box(slide4, Inches(7.1), Inches(5.7), Inches(5), Inches(0.4),
             "Cloud APT Statistics", font_size=12, bold=True, color=LIGHT_ACCENT)
add_text_box(slide4, Inches(7.3), Inches(6.15), Inches(2.2), Inches(0.7),
             "82%\nBreaches involve cloud", font_size=14, bold=True, color=WHITE)
add_text_box(slide4, Inches(9.8), Inches(6.15), Inches(2.5), Inches(0.7),
             "45%\nCloud-native attacks", font_size=14, bold=True, color=WHITE)

add_slide_number(slide4, 4)


# ============================================================
# SLIDE 5: AI/ML Techniques for Detection
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, DARK_BG)
decorate_slide(slide5)

add_tag(slide5, Inches(0.6), Inches(0.4), "METHODOLOGY", GREEN)
add_text_box(slide5, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8),
             "AI/ML Techniques for APT Detection", font_size=28, bold=True, color=WHITE)
add_text_box(slide5, Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
             "Machine Learning and Deep Learning approaches for detecting sophisticated APT behaviors",
             font_size=12, color=MUTED)

# Three technique cards
techniques = [
    ("Supervised Learning", LIGHT_PRIMARY, [
        ("Random Forest", "96.2%"),
        ("XGBoost", "97.1%"),
        ("SVM", "93.8%"),
    ], "Classification of known attack\npatterns from labeled datasets"),
    ("Unsupervised Learning", LIGHT_ACCENT, [
        ("Isolation Forest", "91.5%"),
        ("Autoencoders", "94.3%"),
        ("DBSCAN", "88.7%"),
    ], "Anomaly detection without\nlabeled APT samples"),
    ("Reinforcement Learning", LIGHT_PURPLE, [
        ("DQN Agent", "89.4%"),
        ("Policy Gradient", "91.2%"),
        ("Multi-Agent RL", "93.6%"),
    ], "Adaptive defense through\ncontinuous learning"),
]

for i, (title, color, metrics, desc) in enumerate(techniques):
    x = Inches(0.4 + i * 4.2)
    add_rounded_rect(slide5, x, Inches(2.1), Inches(3.9), Inches(4.5),
                     RGBColor(0x1E, 0x29, 0x3B), color)
    add_text_box(slide5, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.5),
                 title, font_size=14, bold=True, color=color)
    add_text_box(slide5, x + Inches(0.2), Inches(2.8), Inches(3.5), Inches(0.8),
                 desc, font_size=10, color=MUTED)
    # Metrics
    for j, (method, acc) in enumerate(metrics):
        y = Inches(3.7 + j * 0.7)
        add_text_box(slide5, x + Inches(0.2), y, Inches(2.5), Inches(0.3),
                     method, font_size=10, color=MUTED)
        add_text_box(slide5, x + Inches(2.8), y, Inches(0.9), Inches(0.3),
                     acc, font_size=10, bold=True, color=GREEN, alignment=PP_ALIGN.RIGHT)
        # Progress bar background
        bar_width = float(acc.strip('%')) / 100 * 3.2
        add_rounded_rect(slide5, x + Inches(0.2), y + Inches(0.3),
                         Inches(3.5), Inches(0.12),
                         RGBColor(0x1A, 0x1F, 0x2E), None)
        add_rounded_rect(slide5, x + Inches(0.2), y + Inches(0.3),
                         Inches(bar_width), Inches(0.12), color, None)

# Flow diagram at bottom
add_rounded_rect(slide5, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.55),
                 RGBColor(0x1E, 0x29, 0x3B), PRIMARY)
flow_items = ["Network Logs", ">", "Feature Extraction", ">", "AI/ML Model", ">", "Threat Classification", ">", "Alert & Response"]
flow_x = 0.6
for item in flow_items:
    if item == ">":
        add_text_box(slide5, Inches(flow_x), Inches(6.85), Inches(0.4), Inches(0.3),
                     ">", font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)
        flow_x += 0.4
    else:
        add_text_box(slide5, Inches(flow_x), Inches(6.85), Inches(2.0), Inches(0.3),
                     item, font_size=10, bold=True, color=LIGHT_PRIMARY, alignment=PP_ALIGN.CENTER)
        flow_x += 2.2

add_slide_number(slide5, 5)


# ============================================================
# SLIDE 6: Deep Learning Approaches
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, DARK_BG)
decorate_slide(slide6)

add_tag(slide6, Inches(0.6), Inches(0.4), "DEEP LEARNING", LIGHT_PRIMARY)
add_text_box(slide6, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8),
             "Advanced Deep Learning Architectures", font_size=28, bold=True, color=WHITE)
add_text_box(slide6, Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
             "State-of-the-art neural network architectures for APT detection in cloud environments",
             font_size=12, color=MUTED)

dl_methods = [
    ("CNN + LSTM Hybrid", "Convolutional layers extract spatial features from network flow matrices while LSTM captures temporal attack sequences",
     "Spatial Features | Temporal Patterns", "98.3%", LIGHT_PRIMARY),
    ("Graph Neural Networks", "Model cloud network topology as graphs; detect anomalous lateral movement and suspicious communication patterns",
     "Topology-Aware | Relationship Mining", "96.7%", LIGHT_ACCENT),
    ("Transformer Models (LLM)", "Attention mechanisms for long-range dependency modeling in system call sequences and log analysis (APT-LLM)",
     "Self-Attention | Log Analysis", "97.8%", LIGHT_PURPLE),
    ("Generative Adversarial Networks", "GANs generate synthetic APT traffic for training data augmentation; Discriminator as anomaly detector",
     "Data Augmentation | Anomaly Detection", "95.2%", PINK),
]

for i, (title, desc, tags, acc, color) in enumerate(dl_methods):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(2.1 + row * 2.5)
    add_rounded_rect(slide6, x, y, Inches(6.0), Inches(2.2),
                     RGBColor(0x1E, 0x29, 0x3B), color)
    add_text_box(slide6, x + Inches(0.3), y + Inches(0.2), Inches(5.4), Inches(0.4),
                 title, font_size=13, bold=True, color=color)
    add_text_box(slide6, x + Inches(0.3), y + Inches(0.6), Inches(5.4), Inches(0.8),
                 desc, font_size=10, color=MUTED)
    add_text_box(slide6, x + Inches(0.3), y + Inches(1.4), Inches(4.0), Inches(0.3),
                 tags, font_size=9, color=MUTED)
    add_text_box(slide6, x + Inches(4.5), y + Inches(1.4), Inches(1.2), Inches(0.3),
                 acc, font_size=12, bold=True, color=GREEN, alignment=PP_ALIGN.RIGHT)

add_slide_number(slide6, 6)


# ============================================================
# SLIDE 7: Literature Review
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, DARK_BG)
decorate_slide(slide7)

add_tag(slide7, Inches(0.6), Inches(0.4), "LITERATURE REVIEW", LIGHT_ACCENT)
add_text_box(slide7, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8),
             "Key Research Contributions", font_size=28, bold=True, color=WHITE)
add_text_box(slide7, Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
             "Comprehensive review of state-of-the-art approaches in AI-based APT detection",
             font_size=12, color=MUTED)

# Table header
add_rounded_rect(slide7, Inches(0.4), Inches(2.1), Inches(12.5), Inches(0.5),
                 RGBColor(0x1E, 0x1F, 0x3D), PRIMARY)
headers = [("Study / Year", 0.6), ("Technique", 3.8), ("Focus Area", 6.2), ("Key Finding", 9.0)]
for text, x in headers:
    add_text_box(slide7, Inches(x), Inches(2.15), Inches(2.5), Inches(0.4),
                 text, font_size=9, bold=True, color=LIGHT_PRIMARY)

# Table rows
rows = [
    ("CMU SEI Report (2024)", "AI/ML Survey", "APT Defense Feasibility", "Layered AI approach recommended"),
    ("APT-LLM (2025)", "Large Language Models", "Anomaly Detection", "LLM embeddings outperform traditional"),
    ("RAPID (2024)", "Context-Aware DL", "Self-supervised Detection", "Adapts to dynamic system behavior"),
    ("Nature (2024)", "Advanced Computing", "Network Traffic Profiling", "Behavioral APT profiles in traffic"),
    ("DeepTaskAPT (2021)", "LSTM", "Insider Threat", "Task-tree sequence modeling"),
    ("Hybrid Ensemble (2025)", "Ensemble ML", "Network APT Detection", "Optimized hybrid classifiers"),
    ("Cloud-Edge Synergy (2026)", "Graph Transformer", "Collaborative Detection", "Distributed cloud-edge framework"),
]

for i, (study, tech, focus, finding) in enumerate(rows):
    y = Inches(2.7 + i * 0.55)
    bg_color = RGBColor(0x15, 0x1D, 0x2F) if i % 2 == 0 else RGBColor(0x12, 0x1A, 0x2A)
    add_rounded_rect(slide7, Inches(0.4), y, Inches(12.5), Inches(0.5), bg_color, None)
    add_text_box(slide7, Inches(0.6), y + Inches(0.05), Inches(3.0), Inches(0.4),
                 study, font_size=10, bold=True, color=WHITE)
    add_text_box(slide7, Inches(3.8), y + Inches(0.05), Inches(2.2), Inches(0.4),
                 tech, font_size=10, color=MUTED)
    add_text_box(slide7, Inches(6.2), y + Inches(0.05), Inches(2.6), Inches(0.4),
                 focus, font_size=10, color=MUTED)
    add_text_box(slide7, Inches(9.0), y + Inches(0.05), Inches(3.5), Inches(0.4),
                 finding, font_size=10, color=MUTED)

# Bottom stats
stats7 = [("34+", "IDS Models for Cloud"), ("95%+", "Accuracy (Hybrid)"), ("2020-2026", "Research Timeline")]
for i, (val, label) in enumerate(stats7):
    x = Inches(0.8 + i * 4.2)
    add_rounded_rect(slide7, x, Inches(6.6), Inches(3.5), Inches(0.8),
                     RGBColor(0x1E, 0x29, 0x3B), PRIMARY)
    add_text_box(slide7, x, Inches(6.65), Inches(3.5), Inches(0.4),
                 val, font_size=16, bold=True, color=LIGHT_ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, x, Inches(7.0), Inches(3.5), Inches(0.3),
                 label, font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide7, 7)


# ============================================================
# SLIDE 8: Research Gaps
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, DARK_BG)
decorate_slide(slide8)

add_tag(slide8, Inches(0.6), Inches(0.4), "CRITICAL ANALYSIS", RED)
add_text_box(slide8, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8),
             "Identified Research Gaps", font_size=28, bold=True, color=WHITE)
add_text_box(slide8, Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
             "Key limitations in existing literature forming the basis for our proposed research",
             font_size=12, color=MUTED)

gaps = [
    ("1", "Lack of Cloud-Native APT Datasets", "Most studies use generic IDS datasets (NSL-KDD, CICIDS) that don't reflect cloud-specific attack patterns", RED),
    ("2", "Limited Real-Time Processing", "Deep learning models have high inference latency for high-throughput cloud environments", ORANGE),
    ("3", "Explainability Deficit (XAI)", "Black-box AI models lack interpretability for SOC analysts; need for explainable threat attribution", YELLOW),
    ("4", "Multi-Cloud & Hybrid Environments", "Few studies address APT detection across multi-cloud or hybrid cloud-edge architectures", GREEN),
    ("5", "Adversarial Robustness", "ML models vulnerable to adversarial evasion; APT actors can craft traffic to fool detectors", LIGHT_ACCENT),
    ("6", "Cross-Stage Correlation", "Limited work on correlating alerts across multiple kill chain stages for holistic detection", LIGHT_PRIMARY),
    ("7", "Federated & Privacy-Preserving", "Collaborative threat intelligence without sharing sensitive cloud tenant data remains unsolved", LIGHT_PURPLE),
    ("8", "Concept Drift Handling", "APT tactics evolve rapidly; models degrade over time without continuous adaptation", PINK),
]

for i, (num, title, desc, color) in enumerate(gaps):
    col = i // 4
    row = i % 4
    x = Inches(0.4 + col * 6.4)
    y = Inches(2.1 + row * 1.2)
    # Number circle
    add_rounded_rect(slide8, x, y, Inches(0.4), Inches(0.4), RGBColor(0x1E, 0x29, 0x3B), color)
    add_text_box(slide8, x, y + Inches(0.02), Inches(0.4), Inches(0.35),
                 num, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Text
    add_text_box(slide8, x + Inches(0.5), y, Inches(5.5), Inches(0.35),
                 title, font_size=12, bold=True, color=WHITE)
    add_text_box(slide8, x + Inches(0.5), y + Inches(0.35), Inches(5.5), Inches(0.7),
                 desc, font_size=10, color=MUTED)

# Bottom opportunity box
add_rounded_rect(slide8, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.5),
                 RGBColor(0x2D, 0x1B, 0x1B), RED)
add_text_box(slide8, Inches(0.6), Inches(6.93), Inches(12), Inches(0.45),
             "Primary Opportunity: Explainable, real-time, cloud-native AI framework correlating multi-stage APT behaviors with adversarial robustness",
             font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide8, 8)


# ============================================================
# SLIDE 9: Proposed Research Direction
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, DARK_BG)
decorate_slide(slide9)

add_tag(slide9, Inches(0.6), Inches(0.4), "OUR CONTRIBUTION", GREEN)
add_text_box(slide9, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8),
             "Proposed Research Direction", font_size=28, bold=True, color=WHITE)
add_text_box(slide9, Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
             "A novel AI-driven framework addressing identified gaps for next-gen APT detection in cloud networks",
             font_size=12, color=MUTED)

# Left - Research Objectives
add_rounded_rect(slide9, Inches(0.4), Inches(2.1), Inches(6.0), Inches(2.5),
                 RGBColor(0x1E, 0x29, 0x3B), GREEN)
add_text_box(slide9, Inches(0.7), Inches(2.3), Inches(5.5), Inches(0.4),
             "Research Objectives", font_size=13, bold=True, color=GREEN)
objectives = [
    "Design cloud-native AI architecture for real-time APT detection",
    "Implement explainable multi-stage threat correlation",
    "Address adversarial robustness against evasion attacks",
    "Enable federated learning for privacy-preserving detection",
    "Evaluate on realistic cloud APT scenarios",
]
for i, obj in enumerate(objectives):
    add_text_box(slide9, Inches(0.7), Inches(2.8 + i * 0.35), Inches(5.5), Inches(0.35),
                 f"  {obj}", font_size=10, color=MUTED)

# Left - Methodology
add_rounded_rect(slide9, Inches(0.4), Inches(4.8), Inches(6.0), Inches(2.5),
                 RGBColor(0x1E, 0x29, 0x3B), LIGHT_PRIMARY)
add_text_box(slide9, Inches(0.7), Inches(5.0), Inches(5.5), Inches(0.4),
             "Methodology", font_size=13, bold=True, color=LIGHT_PRIMARY)
phases = [
    ("Phase 1:", "Cloud APT dataset generation using MITRE ATT&CK simulation"),
    ("Phase 2:", "Hybrid GNN + Transformer model development"),
    ("Phase 3:", "XAI integration with SHAP/LIME explanations"),
    ("Phase 4:", "Federated deployment across multi-cloud testbed"),
]
for i, (phase, desc) in enumerate(phases):
    add_text_box(slide9, Inches(0.7), Inches(5.5 + i * 0.4), Inches(1.0), Inches(0.35),
                 phase, font_size=10, bold=True, color=WHITE)
    add_text_box(slide9, Inches(1.7), Inches(5.5 + i * 0.4), Inches(4.5), Inches(0.35),
                 desc, font_size=10, color=MUTED)

# Right - Expected Contributions
add_rounded_rect(slide9, Inches(6.8), Inches(2.1), Inches(5.8), Inches(2.8),
                 RGBColor(0x1E, 0x29, 0x3B), LIGHT_ACCENT)
add_text_box(slide9, Inches(7.1), Inches(2.3), Inches(5.2), Inches(0.4),
             "Expected Contributions", font_size=13, bold=True, color=LIGHT_ACCENT)
contribs = [
    ("Novel Framework:", "Cloud-native, explainable APT detection system"),
    ("New Dataset:", "Realistic cloud APT attack dataset for community"),
    ("Benchmark:", "Comparison across 10+ state-of-the-art methods"),
    ("Open Source:", "Reproducible implementation for research"),
]
for i, (key, val) in enumerate(contribs):
    add_text_box(slide9, Inches(7.1), Inches(2.85 + i * 0.5), Inches(1.8), Inches(0.4),
                 key, font_size=10, bold=True, color=WHITE)
    add_text_box(slide9, Inches(8.9), Inches(2.85 + i * 0.5), Inches(3.5), Inches(0.4),
                 val, font_size=10, color=MUTED)

# Right - Innovation tags
add_rounded_rect(slide9, Inches(6.8), Inches(5.2), Inches(5.8), Inches(2.1),
                 RGBColor(0x1E, 0x29, 0x3B), PURPLE)
add_text_box(slide9, Inches(7.1), Inches(5.4), Inches(5.2), Inches(0.4),
             "Innovation Highlights", font_size=13, bold=True, color=LIGHT_PURPLE)
innovations = ["Graph Transformers", "Federated Learning", "Explainable AI",
               "Cloud-Native", "Adversarial Training", "Kill Chain Mapping"]
for i, inn in enumerate(innovations):
    row = i // 3
    col = i % 3
    add_text_box(slide9, Inches(7.1 + col * 1.8), Inches(5.9 + row * 0.5),
                 Inches(1.7), Inches(0.4), inn, font_size=9, bold=True, color=LIGHT_ACCENT)

add_slide_number(slide9, 9)


# ============================================================
# SLIDE 10: References
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, DARK_BG)
decorate_slide(slide10)

add_tag(slide10, Inches(0.6), Inches(0.4), "BIBLIOGRAPHY", LIGHT_PRIMARY)
add_text_box(slide10, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8),
             "Key References & Resources", font_size=28, bold=True, color=WHITE)

references = [
    "[1] CMU SEI (2024) - Toward the Use of AI for APT Detection\n     https://sei.cmu.edu/library/toward-the-use-of-artificial-intelligence-ai-for-advanced-persistent-threat-detection/",
    "[2] Arxiv (2025) - APT-LLM: Embedding-Based Anomaly Detection Using LLMs\n     https://arxiv.org/html/2502.09385v1",
    "[3] Arxiv (2024) - RAPID: Context-Aware Deep Learning for Robust APT Detection\n     https://arxiv.org/abs/2406.05362",
    "[4] Nature Scientific Reports (2024) - Novel Approach for APT Detection\n     https://www.nature.com/articles/s41598-024-72957-0",
    "[5] Springer (2024) - Explainable Deep Learning for APTs: A Review\n     https://link.springer.com/article/10.1007/s10462-024-10890-4",
    "[6] Oxford Academic (2024) - SLR on APT Behaviors and Detection Strategy\n     https://academic.oup.com/cybersecurity/article/10/1/tyad023/7504935",
]

references2 = [
    "[7] ACM (2024) - APT Detection Systems: Review of Approaches & Trends\n     https://dl.acm.org/doi/10.1145/3696014",
    "[8] Journal of Big Data (2025) - Hybrid Ensemble ML for Detecting APTs\n     https://journalofbigdata.springeropen.com/articles/10.1186/s40537-025-01272-w",
    "[9] Springer (2026) - Collaborative APT Detection: Cloud-Edge Graph Transformer\n     https://link.springer.com/article/10.1186/s13677-026-00940-3",
    "[10] Journal of Big Data (2024) - AI-Driven Detection Techniques Review\n      https://journalofbigdata.springeropen.com/articles/10.1186/s40537-024-00957-y",
    "[11] ResearchGate (2024) - APT Attribution Using Deep Reinforcement Learning\n      https://www.researchgate.net/publication/384937397",
    "[12] ResearchGate (2025) - AI-Powered Real-time Threat Detection in Cloud\n      https://www.researchgate.net/publication/391269229",
]

# Left column
for i, ref in enumerate(references):
    y = Inches(1.5 + i * 0.85)
    add_rounded_rect(slide10, Inches(0.4), y, Inches(6.2), Inches(0.78),
                     RGBColor(0x15, 0x1D, 0x2F), PRIMARY)
    add_text_box(slide10, Inches(0.6), y + Inches(0.05), Inches(5.9), Inches(0.7),
                 ref, font_size=8, color=MUTED)

# Right column
for i, ref in enumerate(references2):
    y = Inches(1.5 + i * 0.85)
    add_rounded_rect(slide10, Inches(6.9), y, Inches(6.2), Inches(0.78),
                     RGBColor(0x15, 0x1D, 0x2F), PRIMARY)
    add_text_box(slide10, Inches(7.1), y + Inches(0.05), Inches(5.9), Inches(0.7),
                 ref, font_size=8, color=MUTED)

add_slide_number(slide10, 10)


# ============================================================
# Save the presentation
# ============================================================
output_path = "/projects/sandbox/APT_Detection_Cloud_AI_Research.pptx"
prs.save(output_path)
print(f"Presentation saved successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
