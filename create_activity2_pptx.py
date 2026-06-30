"""
Activity 2: Data Collection Methods - PPTX
Topic: Detection of APT in Cloud Networks Using AI
3 slides with professional dark theme and humanized wording
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Colors
BG = RGBColor(0x0B, 0x11, 0x20)
CARD = RGBColor(0x1A, 0x23, 0x32)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LT = RGBColor(0x93, 0xC5, 0xFD)
TEAL = RGBColor(0x08, 0x91, 0xB2)
TEAL_LT = RGBColor(0x67, 0xE8, 0xF9)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LT = RGBColor(0xC4, 0xB5, 0xFD)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_LT = RGBColor(0x6E, 0xE7, 0xB7)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SEC = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def orb(slide, l, t, w, h, color, alpha=0.87):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    spPr = s._element.find(qn('p:spPr'))
    if spPr is not None:
        sf = spPr.find('.//' + qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha'))
                a.set('val', str(int((1 - alpha) * 100000)))

def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tb

def decor(slide):
    orb(slide, Inches(10.5), Inches(-0.8), Inches(3), Inches(3), TEAL)
    orb(slide, Inches(-0.8), Inches(5.5), Inches(2.2), Inches(2.2), BLUE, 0.9)

def snum(slide, n):
    txt(slide, Inches(11.8), Inches(6.95), Inches(1.2), Inches(0.35),
        f"{n:02d} / 03", sz=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: Introduction to Data Collection
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)
decor(s1)
orb(s1, Inches(5), Inches(2.5), Inches(4), Inches(4), PURPLE, 0.93)

txt(s1, Inches(2.5), Inches(0.6), Inches(8), Inches(0.35),
    "ACTIVITY 2: DATA COLLECTION METHODS",
    sz=9, bold=True, color=TEAL_LT, align=PP_ALIGN.CENTER)
txt(s1, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.0),
    "Selecting Two Data Collection Methods",
    sz=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, Inches(1.8), Inches(2.3), Inches(9.7), Inches(0.9),
    "With a quantitative methodology in place, we now need to decide how we will "
    "actually gather the data our AI models require. The strategy must provide "
    "large-scale, structured, and labelled datasets for rigorous training and testing.",
    sz=12, color=SEC, align=PP_ALIGN.CENTER)

# Left card - Method 1
rect(s1, Inches(0.5), Inches(3.5), Inches(6.0), Inches(3.7), CARD, BLUE)
txt(s1, Inches(0.8), Inches(3.7), Inches(5.4), Inches(0.4),
    "Method 1: Secondary Data Analysis", sz=13, bold=True, color=BLUE_LT)
txt(s1, Inches(0.8), Inches(4.05), Inches(5.4), Inches(0.25),
    "Publicly Available Benchmark Datasets", sz=9, color=MUTED)
txt(s1, Inches(0.8), Inches(4.4), Inches(5.4), Inches(1.5),
    "We draw on established, peer-reviewed network intrusion datasets "
    "that the cybersecurity research community has validated over many years. "
    "These contain both normal traffic and carefully labelled attack samples, "
    "giving us a solid foundation for model training and direct comparison "
    "with published work.",
    sz=10, color=MUTED)
datasets = ["CICIDS 2017/2018", "UNSW-NB15", "Unraveled APT Dataset", "DARPA Transparent Computing"]
for i, ds in enumerate(datasets):
    txt(s1, Inches(1.0), Inches(5.9 + i * 0.35), Inches(5.0), Inches(0.3),
        f"  {ds}", sz=9, color=SEC)

# Right card - Method 2
rect(s1, Inches(6.8), Inches(3.5), Inches(6.0), Inches(3.7), CARD, TEAL)
txt(s1, Inches(7.1), Inches(3.7), Inches(5.4), Inches(0.4),
    "Method 2: Controlled Experiment", sz=13, bold=True, color=TEAL_LT)
txt(s1, Inches(7.1), Inches(4.05), Inches(5.4), Inches(0.25),
    "Simulated Cloud Environment with Live Attacks", sz=9, color=MUTED)
txt(s1, Inches(7.1), Inches(4.4), Inches(5.4), Inches(1.5),
    "We build a cloud-based testbed and run realistic attack campaigns "
    "against it using established adversary simulation frameworks. This "
    "produces fresh, cloud-native data that captures modern threat behaviours "
    "not found in older benchmark datasets.",
    sz=10, color=MUTED)
tools = ["AWS/Azure Cloud Testbed", "MITRE ATT&CK Mapping", "Atomic Red Team Tools", "Custom APT Scenarios"]
for i, t in enumerate(tools):
    txt(s1, Inches(7.3), Inches(5.9 + i * 0.35), Inches(5.0), Inches(0.3),
        f"  {t}", sz=9, color=SEC)

snum(s1, 1)


# ============================================================
# SLIDE 2: Justification for Each Method
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
decor(s2)

txt(s2, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
    "DETAILED JUSTIFICATION", sz=9, bold=True, color=BLUE_LT)
txt(s2, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "Why These Methods Work for Our Study", sz=26, bold=True, color=WHITE)
txt(s2, Inches(0.6), Inches(1.45), Inches(10), Inches(0.4),
    "Each method addresses a different aspect of our data needs, ensuring both breadth and depth.",
    sz=11, color=MUTED)

# Left - Secondary Data justification
rect(s2, Inches(0.4), Inches(2.0), Inches(6.2), Inches(5.2), CARD, BLUE)
txt(s2, Inches(0.7), Inches(2.2), Inches(5.5), Inches(0.4),
    "Secondary Data - The Case For It", sz=13, bold=True, color=BLUE_LT)

m1_reasons = [
    ("Community-validated", "Used in hundreds of published studies, enabling direct comparison with existing work"),
    ("Pre-labelled attacks", "Each record already tagged, eliminating manual annotation and reducing human bias"),
    ("Massive scale", "Millions of network flow records providing enough volume for data-hungry deep learning"),
    ("Reproducibility", "Anyone can download the same data and replicate our experiments exactly"),
]
for i, (title, desc) in enumerate(m1_reasons):
    y = Inches(2.8 + i * 0.7)
    txt(s2, Inches(0.9), y, Inches(5.3), Inches(0.3),
        f"  {title}", sz=10, bold=True, color=WHITE)
    txt(s2, Inches(0.9), y + Inches(0.28), Inches(5.3), Inches(0.35),
        f"    {desc}", sz=9, color=MUTED)

# Method 1 callout
rect(s2, Inches(0.5), Inches(5.7), Inches(5.9), Inches(1.0),
     RGBColor(0x0A, 0x18, 0x2E), BLUE)
txt(s2, Inches(0.7), Inches(5.8), Inches(5.5), Inches(0.8),
    "How it helps: These datasets serve as our training foundation and "
    "allow benchmarking against state-of-the-art methods using identical "
    "evaluation conditions.",
    sz=9, color=SEC)

# Right - Controlled Experiment justification
rect(s2, Inches(6.9), Inches(2.0), Inches(5.8), Inches(5.2), CARD, TEAL)
txt(s2, Inches(7.2), Inches(2.2), Inches(5.2), Inches(0.4),
    "Controlled Experiment - The Case For It", sz=13, bold=True, color=TEAL_LT)

m2_reasons = [
    ("Cloud-native patterns", "Captures API abuse, container escape, and lateral movement unique to cloud"),
    ("Current threats", "Simulates 2024-2025 APT tactics that older datasets simply cannot reflect"),
    ("Full variable control", "We decide exactly which attacks to run, when, and at what intensity"),
    ("Addresses research gap", "Multiple review papers note lack of cloud-specific APT datasets"),
]
for i, (title, desc) in enumerate(m2_reasons):
    y = Inches(2.8 + i * 0.7)
    txt(s2, Inches(7.4), y, Inches(5.0), Inches(0.3),
        f"  {title}", sz=10, bold=True, color=WHITE)
    txt(s2, Inches(7.4), y + Inches(0.28), Inches(5.0), Inches(0.35),
        f"    {desc}", sz=9, color=MUTED)

# Method 2 callout
rect(s2, Inches(7.0), Inches(5.7), Inches(5.5), Inches(1.0),
     RGBColor(0x0A, 0x1C, 0x22), TEAL)
txt(s2, Inches(7.2), Inches(5.8), Inches(5.1), Inches(0.8),
    "How it helps: Validates our model on fresh, realistic cloud attack "
    "scenarios representing the threats organisations actually face today.",
    sz=9, color=SEC)

snum(s2, 2)


# ============================================================
# SLIDE 3: Summary & How Methods Support Research
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
decor(s3)

txt(s3, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
    "BRINGING IT TOGETHER", sz=9, bold=True, color=GREEN_LT)
txt(s3, Inches(0.6), Inches(0.8), Inches(10), Inches(0.6),
    "How Our Data Collection Supports the Research", sz=26, bold=True, color=WHITE)
txt(s3, Inches(0.6), Inches(1.45), Inches(10), Inches(0.4),
    "Both methods work in tandem to ensure robust training and realistic validation.",
    sz=11, color=MUTED)

# Left - Comparison Table
rect(s3, Inches(0.4), Inches(2.0), Inches(6.5), Inches(4.2), CARD, TEAL)
txt(s3, Inches(0.7), Inches(2.2), Inches(6.0), Inches(0.4),
    "Data Collection Overview", sz=13, bold=True, color=TEAL_LT, align=PP_ALIGN.CENTER)

# Table header
rect(s3, Inches(0.5), Inches(2.7), Inches(6.2), Inches(0.4),
     RGBColor(0x12, 0x1A, 0x2A), None)
txt(s3, Inches(0.7), Inches(2.73), Inches(1.5), Inches(0.35),
    "Aspect", sz=8, bold=True, color=TEAL_LT)
txt(s3, Inches(2.3), Inches(2.73), Inches(2.2), Inches(0.35),
    "Method 1: Secondary", sz=8, bold=True, color=BLUE_LT)
txt(s3, Inches(4.6), Inches(2.73), Inches(2.0), Inches(0.35),
    "Method 2: Experiment", sz=8, bold=True, color=TEAL_LT)

table_data = [
    ("Data Source", "Public benchmarks", "Our cloud testbed"),
    ("Scale", "Millions of records", "Thousands (targeted)"),
    ("Labelling", "Pre-labelled by authors", "Ground truth from sim"),
    ("Cloud Specific?", "Partially (generic)", "Fully cloud-native"),
    ("Primary Role", "Training + benchmarking", "Validation + gap filling"),
    ("Reproducible?", "Fully (public access)", "Via documented procedure"),
]
for i, (aspect, m1, m2) in enumerate(table_data):
    y = Inches(3.15 + i * 0.48)
    bg = RGBColor(0x12, 0x1A, 0x2A) if i % 2 == 0 else RGBColor(0x15, 0x1E, 0x2E)
    rect(s3, Inches(0.5), y, Inches(6.2), Inches(0.43), bg, None)
    txt(s3, Inches(0.7), y + Inches(0.04), Inches(1.5), Inches(0.35),
        aspect, sz=9, bold=True, color=WHITE)
    txt(s3, Inches(2.3), y + Inches(0.04), Inches(2.2), Inches(0.35),
        m1, sz=9, color=MUTED)
    txt(s3, Inches(4.6), y + Inches(0.04), Inches(2.0), Inches(0.35),
        m2, sz=9, color=MUTED)

# Right - Why This Combination Works
rect(s3, Inches(7.2), Inches(2.0), Inches(5.6), Inches(3.0), CARD, GREEN)
txt(s3, Inches(7.5), Inches(2.2), Inches(5.0), Inches(0.4),
    "Why This Combination Works", sz=13, bold=True, color=GREEN_LT)

combo_reasons = [
    ("Breadth meets depth", "Benchmarks for fair comparison; experiment for cloud validation"),
    ("Historical meets current", "Established patterns plus latest APT techniques"),
    ("Community meets novelty", "Peers can replicate; simulation contributes new data"),
    ("Aligns with methodology", "Both produce structured numerical data for quantitative analysis"),
]
for i, (title, desc) in enumerate(combo_reasons):
    y = Inches(2.7 + i * 0.6)
    txt(s3, Inches(7.7), y, Inches(5.0), Inches(0.28),
        f"  {title}", sz=10, bold=True, color=WHITE)
    txt(s3, Inches(7.7), y + Inches(0.26), Inches(5.0), Inches(0.3),
        f"    {desc}", sz=9, color=MUTED)

# Bottom final thought
rect(s3, Inches(0.4), Inches(6.4), Inches(12.4), Inches(0.9),
     RGBColor(0x0A, 0x1A, 0x20), GREEN)
txt(s3, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.7),
    "By combining a well-established data source with a purpose-built experimental "
    "environment, we cover both the rigour expected by the academic community and "
    "the practical relevance demanded by industry. Together they provide exactly "
    "the evidence base our research questions require.",
    sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

snum(s3, 3)

# ============================================================
# SAVE
# ============================================================
output = "/projects/sandbox/Activity_2_Data_Collection.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
